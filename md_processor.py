import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches # Use a different alias for docx Inches

def parse_slide_chunk(chunk):
    """Parses a single slide's text block to extract its components."""
    slide_title = "Untitled Slide"
    
    title_match = re.match(r"### (.*)", chunk)
    if title_match:
        slide_title = title_match.group(1).strip()

    text_content_match = re.search(
        r"\*\*Slide Text Content:\*\*(.*?)(?=\*\*Image Prompt:\*\*|\*\*Author Notes for Slide|\Z)", 
        chunk, 
        re.DOTALL
    )
    raw_text_content = text_content_match.group(1).strip() if text_content_match else ""

    image_prompt_match = re.search(
        r"\*\*Image Prompt:\*\*(.*?)(?=\n---\n\*\*Author Notes for Slide|\*\*Author Notes for Slide|\Z)",
        chunk,
        re.DOTALL
    )
    raw_image_prompt = image_prompt_match.group(1).strip() if image_prompt_match else ""
    
    author_notes_match = re.search(
        r"\*\*Author Notes for Slide \d+:\*\*(.*)", 
        chunk, 
        re.DOTALL
    )
    raw_author_notes = author_notes_match.group(1).strip() if author_notes_match else ""
    
    return {
        "title": slide_title,
        "text_content": raw_text_content,
        "image_prompt": raw_image_prompt,
        "author_notes": raw_author_notes
    }

def refined_parse_markdown(markdown_content):
    """Splits the markdown content into slide chunks and parses each chunk."""
    slides_data = []
    processed_slides_content = markdown_content
    match = re.search(r"### Slide \d+:", markdown_content)
    if match:
        processed_slides_content = markdown_content[match.start():]
    
    slide_chunks = re.split(r'\n\s*---\s*\n(?=### Slide \d+:|\Z)', processed_slides_content)

    for chunk in slide_chunks:
        stripped_chunk = chunk.strip()
        if not stripped_chunk or not stripped_chunk.startswith("### Slide"):
            continue
        
        data = parse_slide_chunk(stripped_chunk)
        if data["title"] != "Untitled Slide" or data["text_content"] or data["image_prompt"] or data["author_notes"]:
            slides_data.append(data)
            
    return slides_data

def add_text_to_shape_with_markdown(text_frame, markdown_text, is_title=False):
    """Adds markdown-formatted text to a PowerPoint text frame."""
    text_frame.clear() 
    text_frame.word_wrap = True

    title_font_size = Pt(30)
    content_font_size = Pt(16)

    lines = markdown_text.split('\n')
    if not lines and not markdown_text.strip(): # Handle completely empty text_content
        p = text_frame.add_paragraph()
        run = p.add_run(); run.text = " " 
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = content_font_size if not is_title else title_font_size
        return

    for line_num, line in enumerate(lines):
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        
        # Regex to capture leading spaces (for indentation), bullet char, and content
        bullet_match = re.match(r'^(\s*)([\*\-])\s*(.*)', line) 
        
        is_bullet = False
        content_line = line 
        indent_level = 0

        if bullet_match:
            leading_spaces = bullet_match.group(1)
            # Assuming 2 spaces per indent level for bullets, max level python-pptx supports is 8 (0-indexed)
            indent_level = min(len(leading_spaces) // 2, 8) 
            content_line = bullet_match.group(3) # Actual text after bullet marker and spaces
            is_bullet = True
            p.level = indent_level # Set bullet indentation level
        else:
            # For non-bullet lines, lstrip to remove unintentional leading spaces.
            # Intentional indentation of non-bullet paragraphs is not handled here.
            content_line = line.lstrip() 

        segments = re.split(r'(\*\*.*?\*\*|__.*?__)', content_line)
        
        has_actual_text_in_segments = any(s.strip() for s in segments)

        if not has_actual_text_in_segments: # Content is empty or only whitespace
            if is_bullet:
                # For an empty bullet item (e.g., "* "), add a space to make the bullet visible
                run = p.add_run()
                run.text = " " 
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = content_font_size
            # else: (it's a non-bullet line that's blank or whitespace)
            # The paragraph 'p' is already added. If line.lstrip() was empty, 'p' will be empty.
            # This preserves blank lines.
        else: # Has actual text content
            for segment in segments:
                if not segment: continue 
                
                run = p.add_run()
                if (segment.startswith("**") and segment.endswith("**")) or \
                   (segment.startswith("__") and segment.endswith("__")):
                    run.text = segment[2:-2]
                    run.font.bold = True
                else:
                    run.text = segment
                
                run.font.color.rgb = RGBColor(255, 255, 255)
                if is_title:
                    run.font.size = title_font_size
                else:
                    run.font.size = content_font_size

def create_ppt(slides_data, filename="Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(16) 
    prs.slide_height = Inches(9)

    if not slides_data:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background; fill = background.fill
        fill.solid(); fill.fore_color.rgb = RGBColor(0, 0, 0)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
        tf = title_shape.text_frame; p = tf.add_paragraph(); run = p.add_run()
        run.text = "No slide content found in input."
        run.font.color.rgb = RGBColor(255, 255, 255); run.font.size = Pt(24)
        prs.save(filename)
        print(f"PPTX file '{filename}' created (empty content due to no slide data).")
        return

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[5] 
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background; fill = background.fill
        fill.solid(); fill.fore_color.rgb = RGBColor(0, 0, 0) 

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1.0), Inches(1.0))
        add_text_to_shape_with_markdown(title_shape.text_frame, slide_data["title"], is_title=True)
        
        if slide_data["text_content"]:
            content_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), 
                (prs.slide_width * 0.65) - Inches(0.5), 
                prs.slide_height - Inches(2.0)
            )
            add_text_to_shape_with_markdown(content_shape.text_frame, slide_data["text_content"])

    prs.save(filename)
    print(f"PPTX file '{filename}' created successfully with {len(slides_data)} slides.")

def add_markdown_line_to_docx(doc_parent, markdown_line):
    """Adds a single line of markdown-processed text to a docx document/parent."""
    # Regex to capture leading spaces (for indentation), bullet char, and content
    bullet_match = re.match(r'^(\s*)([\*\-])\s*(.*)', markdown_line)

    if bullet_match:
        leading_spaces = bullet_match.group(1)
        content_line = bullet_match.group(3)
        # Assuming 2 spaces per indent level for bullets
        indent_level = min(len(leading_spaces) // 2, 8) 

        p = doc_parent.add_paragraph(style='ListBullet')
        if indent_level > 0:
            # Adjust indentation; 0.25 inches per level is a starting point
            p.paragraph_format.left_indent = DocxInches(0.25 * indent_level) 

        segments = re.split(r'(\*\*.*?\*\*|__.*?__)', content_line)
        added_text_run = False
        for segment in segments:
            if not segment: continue
            run = p.add_run()
            if (segment.startswith("**") and segment.endswith("**")) or \
               (segment.startswith("__") and segment.endswith("__")):
                run.text = segment[2:-2]
                run.font.bold = True
            else:
                run.text = segment
            if segment.strip():
                added_text_run = True
        
        if not added_text_run and not content_line.strip(): # If bullet item is empty (e.g. "* ")
            run = p.add_run()
            run.text = " " # Add a space to make the bullet visible

    else: # Not a bullet point
        stripped_line = markdown_line.lstrip()
        
        if not stripped_line and markdown_line == "": # An actual empty line from markdown
            doc_parent.add_paragraph("") # Preserve blank line
            return
        # If line was only whitespace, lstrip makes it empty, effectively a blank line too
        elif not stripped_line:
            doc_parent.add_paragraph("")
            return

        p = doc_parent.add_paragraph()
        segments = re.split(r'(\*\*.*?\*\*|__.*?__)', stripped_line)
        for segment in segments:
            if not segment: continue
            run = p.add_run()
            if (segment.startswith("**") and segment.endswith("**")) or \
               (segment.startswith("__") and segment.endswith("__")):
                run.text = segment[2:-2]
                run.font.bold = True
            else:
                run.text = segment
        
        # If paragraph is empty after processing (e.g. line was '****' or just whitespace)
        # and the original stripped_line was also empty, it's fine (empty paragraph).
        # If stripped_line was not empty but p.runs is empty, it means it was like '****'.
        # An empty run is already added by the loop.

def create_doc(slides_data, doc_filename, content_key):
    """Creates a Word document from parsed slide data for a specific content key."""
    doc = Document()
    doc_title = doc_filename.split('.')[0]
    doc.add_heading(doc_title, level=1)
    
    if not slides_data:
        doc.add_paragraph(f"[No slide data to extract for {content_key}]")
        doc.save(doc_filename)
        print(f"DOC file '{doc_filename}' created (empty content due to no slide data).")
        return

    for i, slide_data in enumerate(slides_data):
        doc.add_heading(slide_data["title"], level=2)
        
        content_to_add = slide_data.get(content_key, "") # Get raw content

        if not content_to_add.strip():
            doc.add_paragraph("[No content provided for this section in this slide]")
        else:
            lines = content_to_add.split('\n')
            for line_text in lines:
                add_markdown_line_to_docx(doc, line_text) 

        if i < len(slides_data) - 1:
            # Add a visual separator, not a markdown '---'
            doc.add_paragraph("\n" + "_" * 30 + "\n")


    doc.save(doc_filename)
    print(f"DOC file '{doc_filename}' created successfully.")

def main():
    input_file = "input.txt"
    try:
        with open(input_file, "r", encoding="utf-8") as f:
            markdown_content = f.read()
    except FileNotFoundError:
        print(f"Error: Input file '{input_file}' not found.")
        print(f"Creating a dummy '{input_file}' for testing purposes with sample content.")
        dummy_content = """### Slide 1: Title Slide (Session 2)

**Slide Text Content:**

*   **A Practical Guide to Quantum Machine Learning**
*   Session 2 of [Total Number of Sessions, e.g., 2]
    *   Indented item 1
    *   Indented item 2
*   [Your Name]

**Image Prompt:**
"A dynamic, modern background for 'A Practical Guide to Quantum Machine Learning'."

---

**Author Notes for Slide 1:**

*   **Topic:** Title for Session 2.
    *   **Sub-topic:** More details.
*   **Spoken Content Guidance:**
    *   "Welcome back everyone!"
    *   "This is **important**."

---
### Slide 2: Recap of Session 1

**Slide Text Content:**

*   **Quick Refresher:**
    *   **Quantum Basics:** Qubits, Gates, Circuits.
        *   Even deeper indent.
    *   **NISQ Era:** Noisy hardware.
*   Another top-level bullet.

**Image Prompt:**
"A collage of key mini-icons from Session 1."

---

**Author Notes for Slide 2:**

*   **Topic:** Recap Session 1.
*   **Spoken Content Guidance:**
    *   "Let's quickly refresh some concepts."
    *   "__Another__ important point."
    *   
    *   A line after a blank line.
"""
        with open(input_file, "w", encoding="utf-8") as f:
            f.write(dummy_content)
        markdown_content = dummy_content

    slides_data = refined_parse_markdown(markdown_content)

    if not slides_data:
        print("No slide data could be parsed from the input file. Output files will reflect this.")
        
    create_ppt(slides_data, filename="Presentation.pptx")
    create_doc(slides_data, "ImagePrompt.DOC", "image_prompt") # Image prompts usually don't have rich markdown
    create_doc(slides_data, "OtherNotes.DOC", "author_notes")

if __name__ == "__main__":
    main()
